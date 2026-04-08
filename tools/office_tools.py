"""
Aura-X — Office Document Tools
Creates Word, Excel, and PowerPoint files using python-docx, openpyxl, python-pptx.
Falls back to plain-text if those packages aren't installed.
"""

import os
from pathlib import Path
from typing import Dict
from core.logger import setup_logger

logger = setup_logger("aura_x.tools.office")


def handle_create_word(params: Dict) -> Dict:
    """Create a Word document."""
    filepath = params.get("filepath") or params.get("filename", "document.docx")
    content = params.get("content", "")
    title = params.get("title", "")
    output_dir = params.get("output_dir", str(Path.home() / "Documents"))

    # Build full path if only filename given
    fp = Path(filepath)
    if not fp.is_absolute():
        fp = Path(output_dir) / fp
    fp.parent.mkdir(parents=True, exist_ok=True)

    # Ensure .docx extension
    if fp.suffix.lower() != ".docx":
        fp = fp.with_suffix(".docx")

    try:
        from docx import Document
        doc = Document()
        if title:
            doc.add_heading(title, level=0)
        if content:
            for para in content.split("\n"):
                doc.add_paragraph(para)
        doc.save(str(fp))
        return {"status": "success", "message": f"Word document created: {fp}", "filepath": str(fp)}
    except ImportError:
        # Fallback: save as .txt
        txt_path = fp.with_suffix(".txt")
        text = f"{title}\n{'='*len(title)}\n\n{content}" if title else content
        txt_path.write_text(text, encoding="utf-8")
        return {"status": "success",
                "message": f"python-docx not installed. Saved as text: {txt_path}",
                "filepath": str(txt_path)}
    except Exception as e:
        logger.error(f"Word creation error: {e}", exc_info=True)
        return {"status": "error", "error": str(e)}


def handle_create_excel(params: Dict) -> Dict:
    """Create an Excel spreadsheet."""
    filepath = params.get("filepath") or params.get("filename", "spreadsheet.xlsx")
    headers = params.get("headers", [])
    data = params.get("data", [])
    title = params.get("title", "Sheet1")
    output_dir = params.get("output_dir", str(Path.home() / "Documents"))

    fp = Path(filepath)
    if not fp.is_absolute():
        fp = Path(output_dir) / fp
    fp.parent.mkdir(parents=True, exist_ok=True)
    if fp.suffix.lower() != ".xlsx":
        fp = fp.with_suffix(".xlsx")

    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = title
        if headers:
            ws.append(headers)
        for row in data:
            ws.append(row if isinstance(row, list) else [row])
        wb.save(str(fp))
        return {"status": "success", "message": f"Excel created: {fp}", "filepath": str(fp)}
    except ImportError:
        # Fallback: CSV
        csv_path = fp.with_suffix(".csv")
        lines = []
        if headers:
            lines.append(",".join(str(h) for h in headers))
        for row in data:
            if isinstance(row, list):
                lines.append(",".join(str(c) for c in row))
            else:
                lines.append(str(row))
        csv_path.write_text("\n".join(lines), encoding="utf-8")
        return {"status": "success",
                "message": f"openpyxl not installed. Saved as CSV: {csv_path}",
                "filepath": str(csv_path)}
    except Exception as e:
        logger.error(f"Excel creation error: {e}", exc_info=True)
        return {"status": "error", "error": str(e)}


def handle_create_ppt(params: Dict) -> Dict:
    """Create a PowerPoint presentation."""
    filepath = params.get("filepath") or params.get("filename", "presentation.pptx")
    title = params.get("title", "Presentation")
    slides = params.get("slides", [])
    output_dir = params.get("output_dir", str(Path.home() / "Documents"))

    fp = Path(filepath)
    if not fp.is_absolute():
        fp = Path(output_dir) / fp
    fp.parent.mkdir(parents=True, exist_ok=True)
    if fp.suffix.lower() != ".pptx":
        fp = fp.with_suffix(".pptx")

    try:
        from pptx import Presentation
        prs = Presentation()

        # Title slide
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        # Content slides
        for s in slides:
            layout = prs.slide_layouts[1]
            sl = prs.slides.add_slide(layout)
            if isinstance(s, dict):
                sl.shapes.title.text = s.get("title", "")
                if s.get("content"):
                    sl.placeholders[1].text = s["content"]
            elif isinstance(s, str):
                sl.shapes.title.text = s

        prs.save(str(fp))
        return {"status": "success", "message": f"PowerPoint created: {fp}", "filepath": str(fp)}
    except ImportError:
        # Fallback: text outline
        txt_path = fp.with_suffix(".txt")
        lines = [f"# {title}", ""]
        for i, s in enumerate(slides, 1):
            if isinstance(s, dict):
                lines.append(f"## Slide {i}: {s.get('title', '')}")
                lines.append(s.get("content", ""))
            else:
                lines.append(f"## Slide {i}: {s}")
            lines.append("")
        txt_path.write_text("\n".join(lines), encoding="utf-8")
        return {"status": "success",
                "message": f"python-pptx not installed. Saved as text: {txt_path}",
                "filepath": str(txt_path)}
    except Exception as e:
        logger.error(f"PPT creation error: {e}", exc_info=True)
        return {"status": "error", "error": str(e)}
